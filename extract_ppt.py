from pptx import Presentation

def extract_text_from_ppt(filepath):
    prs = Presentation(filepath)
    print(f"Total Slides: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(text.encode('ascii', 'ignore').decode('ascii'))
        print("\n")

if __name__ == "__main__":
    extract_text_from_ppt("Hands-On-Machine-Learning-Workshop.pptx")
